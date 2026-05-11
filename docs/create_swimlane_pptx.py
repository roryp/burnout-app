"""Clean swimlane PPTX — labels only, all detail in speaker notes.

Reflects the current two-phase reshape architecture:
  Lane 1 (DETERMINISTIC METRICS): Ingestion → Chaos → Classify+Comply → Stress (BEFORE)
  Lane 2 (RESHAPE): Pre-pass (triageUrgent + defuseChaosInputs) → Supervisor + 6 sub-agents → Mutations
  Lane 3 (OUTPUT): Apply Mutations → Recalculate → Friday Score → Persist
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x47, 0x55, 0x69)
LANE_BG = RGBColor(0x14, 0x1E, 0x30)
LANE_BORDER = RGBColor(0x33, 0x44, 0x55)

C_INGEST = RGBColor(0x3B, 0x82, 0xF6)
C_CHAOS  = RGBColor(0xF5, 0x9E, 0x0B)
C_CLASS  = RGBColor(0xA7, 0x8B, 0xFA)
C_STRESS = RGBColor(0x2D, 0xD4, 0xBF)
C_PREPASS = RGBColor(0x2D, 0xD4, 0xBF)   # teal — same family as deterministic
C_AGENT  = RGBColor(0xFB, 0x71, 0x85)
C_MUTATE = RGBColor(0xEC, 0x4C, 0x6F)
C_OUTPUT = RGBColor(0x81, 0x8C, 0xF8)
C_LOOP   = RGBColor(0xF5, 0x9E, 0x0B)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def lane_bg(slide, l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = LANE_BG
    s.line.color.rgb = LANE_BORDER
    s.line.width = Pt(0.75)
    s.shadow.inherit = False


def rbox(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(c[0] // 5, c[1] // 5, c[2] // 5)
    s.line.color.rgb = c
    s.line.width = Pt(1.5)
    s.shadow.inherit = False
    s.adjustments[0] = 0.15


def txt(slide, l, t, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)


def harrow(slide, x, y, w, color=DIM):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False


def darrow(slide, x, y, h, color=DIM):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.2), h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False


def main():
    prs = Presentation()
    prs.slide_width = SW = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    txt(slide, 0, Inches(0.3), SW, Inches(0.5), "Reshape Pipeline", 34, WHITE, True)

    # Lane geometry
    lx = Inches(2.3); lw = Inches(10.5)
    y1 = Inches(1.1); h1 = Inches(1.5)
    y2 = Inches(2.85); h2 = Inches(1.85)
    y3 = Inches(4.95); h3 = Inches(1.3)

    for y, h in [(y1, h1), (y2, h2), (y3, h3)]:
        lane_bg(slide, lx, y, lw, h)

    # Lane labels
    lbl_x = Inches(0.2); lbl_w = Inches(1.9)
    txt(slide, lbl_x, y1 + Inches(0.45), lbl_w, Inches(0.3), "DETERMINISTIC", 14, C_STRESS, True)
    txt(slide, lbl_x, y1 + Inches(0.75), lbl_w, Inches(0.25), "(metrics)", 11, MUTED)

    txt(slide, lbl_x, y2 + Inches(0.55), lbl_w, Inches(0.3), "RESHAPE", 14, C_AGENT, True)
    txt(slide, lbl_x, y2 + Inches(0.85), lbl_w, Inches(0.25), "(pre-pass + AI)", 11, MUTED)

    txt(slide, lbl_x, y3 + Inches(0.4), lbl_w, Inches(0.3), "OUTPUT", 14, C_OUTPUT, True)

    # ── Lane 1 boxes ──
    bw = Inches(2.15); bh = Inches(0.85)
    by = y1 + Inches(0.32); gap = Inches(0.45)
    x1 = lx + Inches(0.3)
    x2 = x1 + bw + gap
    x3 = x2 + bw + gap
    x4 = x3 + bw + gap

    for x, c, label in [
        (x1, C_INGEST, "Ingestion"),
        (x2, C_CHAOS,  "Chaos Metrics"),
        (x3, C_CLASS,  "Classify + Comply"),
        (x4, C_STRESS, "Stress Score"),
    ]:
        rbox(slide, x, by, bw, bh, c)
        txt(slide, x, by + Inches(0.22), bw, Inches(0.4), label, 17, WHITE, True)

    for x in [x1, x2, x3]:
        harrow(slide, x + bw + Inches(0.05), by + bh / 2 - Inches(0.1), gap - Inches(0.1))

    # Arrow down into Lane 2
    darrow(slide, x4 + bw / 2 - Inches(0.1), by + bh + Inches(0.05),
           y2 - by - bh + Inches(0.25))

    # ── Lane 2 boxes (NEW: two-phase reshape) ──
    # Phase 1: deterministic pre-pass (teal, two stacked tool calls)
    # Phase 2: supervisor + 6 sub-agents (red)
    # Phase 3: label mutations (red-pink)
    p1w = Inches(2.9); p1h = Inches(1.35)
    p2w = Inches(3.6); p2h = Inches(1.35)
    p3w = Inches(2.4); p3h = Inches(1.35)
    pgap = Inches(0.3)

    total_w = p1w + p2w + p3w + 2 * pgap
    px1 = lx + (lw - total_w) / 2
    px2 = px1 + p1w + pgap
    px3 = px2 + p2w + pgap
    py = y2 + Inches(0.3)

    # Phase 1 — deterministic pre-pass
    rbox(slide, px1, py, p1w, p1h, C_PREPASS)
    txt(slide, px1, py + Inches(0.10), p1w, Inches(0.30),
        "DETERMINISTIC PRE-PASS", 11, C_PREPASS, True)
    txt(slide, px1, py + Inches(0.42), p1w, Inches(0.35),
        "triageUrgent(n)", 14, WHITE, True)
    txt(slide, px1, py + Inches(0.75), p1w, Inches(0.35),
        "defuseChaosInputs(clock)", 14, WHITE, True)
    txt(slide, px1, py + Inches(1.05), p1w, Inches(0.25),
        "no LLM — guarantees drop", 10, MUTED)

    # Phase 2 — supervisor + 6 sub-agents
    rbox(slide, px2, py, p2w, p2h, C_AGENT)
    txt(slide, px2, py + Inches(0.10), p2w, Inches(0.30),
        "LANGCHAIN4J SUPERVISOR", 11, C_AGENT, True)
    txt(slide, px2, py + Inches(0.42), p2w, Inches(0.35),
        "6 Sub-Agents", 15, WHITE, True)
    txt(slide, px2, py + Inches(0.74), p2w, Inches(0.30),
        "Triage · Defer · Delegate", 11, MUTED)
    txt(slide, px2, py + Inches(0.93), p2w, Inches(0.30),
        "Classify · Scope · Wellness", 11, MUTED)
    txt(slide, px2, py + Inches(1.10), p2w, Inches(0.22),
        "max 15 invocations · SUMMARY", 9, MUTED)

    # Phase 3 — label mutations
    rbox(slide, px3, py, p3w, p3h, C_MUTATE)
    txt(slide, px3, py + Inches(0.10), p3w, Inches(0.30),
        "@Tool METHODS", 11, C_MUTATE, True)
    txt(slide, px3, py + Inches(0.42), p3w, Inches(0.35),
        "Label Mutations", 14, WHITE, True)
    txt(slide, px3, py + Inches(0.74), p3w, Inches(0.30),
        "AddLabels · RemoveLabels", 10, MUTED)
    txt(slide, px3, py + Inches(0.92), p3w, Inches(0.30),
        "Comment · Unassign", 10, MUTED)
    txt(slide, px3, py + Inches(1.08), p3w, Inches(0.25),
        "SetBody · SetUpdatedAt (legacy)", 10, MUTED)

    # Arrows between phases in Lane 2
    arrow_y = py + p1h / 2 - Inches(0.1)
    harrow(slide, px1 + p1w + Inches(0.02), arrow_y, pgap - Inches(0.04))
    harrow(slide, px2 + p2w + Inches(0.02), arrow_y, pgap - Inches(0.04))

    # Arrow down into Lane 3 (centered on phase 3)
    darrow(slide, px3 + p3w / 2 - Inches(0.1), py + p3h + Inches(0.02),
           y3 - py - p3h + Inches(0.18))

    # ── Lane 3 boxes ──
    obw = Inches(2.2); obh = Inches(0.85)
    oy = y3 + Inches(0.22); ogap = Inches(0.38)
    ox1 = lx + Inches(0.3)
    ox2 = ox1 + obw + ogap
    ox3 = ox2 + obw + ogap
    ox4 = ox3 + obw + ogap

    for ox, c, label in [
        (ox1, C_OUTPUT, "Apply Mutations"),
        (ox2, C_LOOP,   "Recalculate ↻"),
        (ox3, C_STRESS, "Friday Score"),
        (ox4, C_OUTPUT, "Persist Snapshot"),
    ]:
        rbox(slide, ox, oy, obw, obh, c)
        txt(slide, ox, oy + Inches(0.22), obw, Inches(0.4), label, 16, WHITE, True)

    for ox in [ox1, ox2, ox3]:
        harrow(slide, ox + obw + Inches(0.02), oy + obh / 2 - Inches(0.1), ogap - Inches(0.04))

    # Tagline
    txt(slide, 0, Inches(6.45), SW, Inches(0.3),
        "Deterministic metrics  →  Pre-pass guarantees the drop  →  AI rebalances  →  Recalculate  →  Persist",
        12, DIM)

    # ── Speaker notes ──
    slide.notes_slide.notes_text_frame.text = (
        "DETERMINISTIC LANE — metrics calculated on every checkin / flamegraph / reshape\n"
        "1. Ingestion — Issues loaded from in-memory IssueCache (populated by /demo/api/seed, "
        "/demo/api/sync, or MCP sync_issues).\n"
        "2. Chaos Metrics — ChaosMetricsService runs FIRST. Five binary signals each worth 2 points: "
        "mystery meat ≥ 3, urgent ≥ 3, context switching ≥ 6, after-hours, label sprawl ≥ 12. Score 0–10.\n"
        "3. Classify + Comply — IssueClassifierService sorts issues into "
        "DEEP_WORK / QUICK_WIN / MAINTENANCE / DEFERRED (first-match-wins). "
        "ComplianceService checks 8 violation rules (score 100 → 0). "
        "3-3-3 day plan built here.\n"
        "4. Stress Score — WorldState aggregates chaos + classification into 12 capped variables. "
        "Six components summed: Workload (0–40), Chaos (0–30), Context Switching (0–15), "
        "Clarity (0–10), Sustained (0–15), After Hours (0–10). Total capped at 100. "
        "This is the BEFORE score.\n\n"
        "RESHAPE LANE — only runs on POST /demo/api/reshape (and MCP reshape_day)\n"
        "5a. DETERMINISTIC PRE-PASS — runs in BurnoutSupervisorService BEFORE any LLM call:\n"
        "    • triageUrgent(n) is invoked directly for every unassigned-urgent issue, stripping "
        "      `urgent` / `priority:critical` / `priority:high` and adding `triaged,backlog`.\n"
        "    • defuseChaosInputs(clock) replaces empty bodies with a scope-pending placeholder. "
        "      It does NOT rewrite real after-hours / touched-today `updatedAt` values — those "
        "      are real signals about real human activity and are deliberately preserved "
        "      (acknowledge-don't-erase) so the AFTER score and the WellnessAgent gates stay honest.\n"
        "    These two calls guarantee the chaos score drops on every reshape, regardless of which "
        "    sub-agents the LLM picks. The supervisor is told the unassigned-urgent issues are "
        "    already triaged and to leave them alone.\n"
        "5b. LANGCHAIN4J SUPERVISOR — AgenticServices.supervisorBuilder() with maxAgentsInvocations=15 "
        "    and SupervisorResponseStrategy.SUMMARY. Six sub-agents are registered:\n"
        "    • TriageAgent (final cleanup if pre-pass missed anything)\n"
        "    • DeferAgent (defer + next-sprint, unassign)\n"
        "    • DelegateAgent (delegate + needs-owner, unassign)\n"
        "    • ClassifyAgent (quick-win / deep-work / maintenance labels for the 3-3-3 plan)\n"
        "    • ScopeAgent (needs-scope + blocked)\n"
        "    • WellnessAgent (suggestBreak, slowIntake, blockCalendarTime)\n"
        "    Every agent has a deterministic fallback if the LLM is unavailable.\n"
        "5c. LABEL MUTATIONS — Each @Tool method on BurnoutMutationTool emits records from the "
        "    sealed GitHubAction interface: AddLabels, RemoveLabels, Comment, Unassign, SetBody. "
        "    SetUpdatedAt is retained on the sealed interface for backward compatibility but is "
        "    no longer emitted by any code path. Mutations are buffered, not applied directly to GitHub.\n\n"
        "OUTPUT LANE\n"
        "6a. Apply Mutations — DemoFlamegraphController.applyMutationsToIssues writes the buffered "
        "    label/body/timestamp changes to the in-memory IssueCache (not to GitHub).\n"
        "6b. Recalculate — The ENTIRE deterministic pipeline reruns on the mutated issues. "
        "    Chaos drops (urgent stripped, bodies filled), classification re-buckets with the "
        "    new labels, new stress score is computed. This is the AFTER score "
        "    (HIGH → MODERATE in the demo — the after-hours penalty is preserved by design).\n"
        "6c. Friday Score — Deploy readiness: READY ≥ 80, CAUTION ≥ 50, NOT_READY < 50.\n"
        "6d. Persist Snapshot — JPA StressSnapshot with all 6 components + optional self-report. "
        "    Powers the Study Dashboard trend charts."
    )

    out = "docs/Burnout-Swimlanes.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
