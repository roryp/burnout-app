"""
Generate a widescreen (16:9) PowerPoint with the 6 Theoretical Foundations.
Each foundation gets its own slide matching the original card style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
IMG_DIR = os.path.join(os.path.dirname(__file__), "images")
OUT = os.path.join(os.path.dirname(__file__), "Theoretical-Foundations.pptx")

# Colors
DARK_BG = RGBColor(0x0A, 0x3D, 0x62)
HEADER_BG = RGBColor(0x1A, 0x52, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
BODY_TEXT = RGBColor(0x33, 0x33, 0x33)
RED_LABEL = RGBColor(0xC0, 0x39, 0x2B)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE_BG = RGBColor(0xE8, 0xEF, 0xF5)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             color=DARK_TEXT, bold=False, italic=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = "Segoe UI"
    p.alignment = align
    return txBox


def add_bullets(slide, left, top, width, height, items, size=20, color=BODY_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
        p.level = 0
        # Bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        # Remove existing bullets first
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        for existing in pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        pPr.append(buChar)
    return txBox


# ─── Foundation data ─────────────────────────────────────────────────
foundations = [
    {
        "title": "Maslach Burnout Inventory",
        "icon": "\U0001F50B",  # battery
        "key_concept": "3 dimensions: exhaustion, depersonalization, reduced accomplishment",
        "application": [
            "After-hours signals exhaustion",
            "Mystery meat signals depersonalization",
            "No deep work signals reduced accomplishment",
        ],
    },
    {
        "title": "Cognitive Load Theory",
        "icon": "\u2699\uFE0F",  # gear
        "key_concept": "Working memory is limited; extraneous load must be minimized",
        "application": [
            "3-3-3 cap at 7 items",
            "Classification removes ambiguity load",
            "Mystery meat penalty reduces unclear scope",
        ],
    },
    {
        "title": "Yerkes-Dodson Law",
        "icon": "\U0001F4CA",  # bar chart
        "key_concept": "Performance peaks at moderate stress, drops at extremes",
        "application": [
            "Stress score targets 30\u201350 (MODERATE) as optimal zone",
            "Too little stress signals disengagement",
            "Too much stress triggers burnout trajectory",
        ],
    },
    {
        "title": "Deep Work (Cal Newport)",
        "icon": "\U0001F4BC",  # briefcase
        "key_concept": "Sustained focused work requires 90+ min uninterrupted blocks",
        "application": [
            "Calendar fragmentation check",
            "Exactly 1 deep work item per day",
            "Each context switch costs 23 min to resume focus",
        ],
    },
    {
        "title": "Pomodoro / Time Boxing",
        "icon": "\U0001F345",  # tomato
        "key_concept": "Short focused intervals with breaks maintain energy",
        "application": [
            "Quick wins as natural break-points between deep work sessions",
            "Completion of small tasks provides dopamine momentum",
            "Alternating cognitive modes prevents single-channel fatigue",
        ],
    },
    {
        "title": "Plutchik Wheel of Emotions",
        "icon": "\U0001F338",  # cherry blossom
        "key_concept": "8 primary emotions with behavioral signatures",
        "application": [
            "4 emotions detected from GitHub signals:",
            "Frustration \u2022 Exhaustion \u2022 Overwhelm \u2022 Anxiety",
        ],
    },
]


# ═══════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ─── Title slide ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Gradient header bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.2))
bar.fill.solid()
bar.fill.fore_color.rgb = HEADER_BG
bar.line.fill.background()

add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2),
         "Theoretical Foundations", size=48, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(1.4), Inches(10), Inches(0.6),
         "6 psychological frameworks powering Burnout-as-a-Service",
         size=22, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

# 6 mini cards in 3x2 grid
card_w = Inches(3.5)
card_h = Inches(2.2)
gap_x = Inches(0.5)
gap_y = Inches(0.4)
start_x = Inches(0.9)
start_y = Inches(2.8)

for idx, f in enumerate(foundations):
    col = idx % 3
    row = idx // 3
    cx = start_x + col * (card_w + gap_x)
    cy = start_y + row * (card_h + gap_y)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    card.line.fill.background()

    add_text(slide, cx + Inches(0.2), cy + Inches(0.15), card_w - Inches(0.4), Inches(0.5),
             f["icon"], size=32, color=WHITE)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.65), card_w - Inches(0.4), Inches(0.5),
             f["title"], size=16, color=WHITE, bold=True)
    add_text(slide, cx + Inches(0.2), cy + Inches(1.15), card_w - Inches(0.4), Inches(1.0),
             f["key_concept"], size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

# ─── Individual foundation slides ───────────────────────────────────
for idx, f in enumerate(foundations):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SUBTLE_BG)

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()

    # Icon
    add_text(slide, Inches(0.8), Inches(0.25), Inches(1.2), Inches(1.2),
             f["icon"], size=56, color=WHITE)

    # Title
    add_text(slide, Inches(2.2), Inches(0.35), Inches(9), Inches(0.8),
             f["title"], size=40, color=WHITE, bold=True)

    # Subtitle line
    add_text(slide, Inches(2.2), Inches(1.1), Inches(9), Inches(0.5),
             f"Foundation {idx + 1} of 6", size=16,
             color=RGBColor(0xAA, 0xCC, 0xDD))

    # White content card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()

    # Key Concept label
    add_text(slide, Inches(1.3), Inches(2.5), Inches(4), Inches(0.5),
             "Key Concept:", size=20, color=RED_LABEL, bold=True, italic=True)

    # Key Concept bullet
    add_bullets(slide, Inches(1.3), Inches(3.0), Inches(10), Inches(0.8),
                [f["key_concept"]], size=22, color=BODY_TEXT)

    # Application label
    add_text(slide, Inches(1.3), Inches(3.9), Inches(4), Inches(0.5),
             "Application:", size=20, color=RED_LABEL, bold=True, italic=True)

    # Application bullets
    add_bullets(slide, Inches(1.3), Inches(4.4), Inches(10), Inches(2.2),
                f["application"], size=22, color=BODY_TEXT)

    # Try to add foundation image on the right if it exists
    img_name = f"foundation-{idx + 1}-{'maslach' if idx == 0 else 'cognitive-load' if idx == 1 else 'yerkes-dodson' if idx == 2 else 'deep-work' if idx == 3 else 'pomodoro' if idx == 4 else 'plutchik'}.png"
    img_file = os.path.join(IMG_DIR, img_name)
    # We skip embedding the image since the text content IS the image content


# ─── Save ────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
