"""Generate a single-slide PPTX with a horizontal 7-phase pipeline diagram.

Reflects the current two-phase reshape architecture:
  1 Ingestion → 2 Chaos → 3 Classify → 4 Stress (BEFORE)
  → 5 Pre-pass (deterministic) → 6 Supervisor + 6 sub-agents
  → 7 Output (apply + recalc + persist)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ARROW_CLR = RGBColor(0x47, 0x55, 0x69)

PHASES = [
    {"num": "1", "title": "Ingestion",        "color": RGBColor(0x3B, 0x82, 0xF6), "sub": "GitHub Issues\n→ Cache"},
    {"num": "2", "title": "Chaos Metrics",    "color": RGBColor(0xF5, 0x9E, 0x0B), "sub": "5 Signals\n→ Score 0–10"},
    {"num": "3", "title": "Classification",   "color": RGBColor(0xA7, 0x8B, 0xFA), "sub": "4 Buckets\n+ Compliance"},
    {"num": "4", "title": "Stress Score",     "color": RGBColor(0x2D, 0xD4, 0xBF), "sub": "WorldState\n→ 0–100"},
    {"num": "5", "title": "Pre-pass",         "color": RGBColor(0x2D, 0xD4, 0xBF), "sub": "triageUrgent\n+ defuseChaos\n(no LLM)"},
    {"num": "6", "title": "Supervisor + 6",   "color": RGBColor(0xFB, 0x71, 0x85), "sub": "LangChain4j\n6 Sub-Agents\n+ Fallback"},
    {"num": "7", "title": "Output",           "color": RGBColor(0x81, 0x8C, 0xF8), "sub": "Apply · Recalc\n+ Persist"},
]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_box(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Reduce corner rounding
    shape.adjustments[0] = 0.1
    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ARROW_CLR
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.rotation = 0.0
    return shape


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG)

    # Title
    tx = slide.shapes.add_textbox(Inches(0), Inches(0.5), prs.slide_width, Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Burnout-as-a-Service — Algorithm Pipeline"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(0), Inches(1.15), prs.slide_width, Inches(0.4))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = ("Deterministic services compute every metric  •  "
               "Pre-pass guarantees the drop  •  AI agents only rebalance, never decide alone")
    p2.font.size = Pt(13)
    p2.font.color.rgb = MUTED
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER

    # Layout constants (7 phases)
    n = len(PHASES)
    box_w = Inches(1.45)
    box_h = Inches(3.2)
    arrow_w = Inches(0.4)
    arrow_h = Inches(0.3)
    total_w = n * box_w + (n - 1) * arrow_w
    start_x = (prs.slide_width - total_w) / 2
    box_y = Inches(2.4)

    # Light backdrop strips behind phases to visually group:
    # 1–4 = deterministic metrics, 5 = pre-pass (deterministic), 6 = AI, 7 = output
    # Group backdrops are subtle — use very dark tints of lane color.

    def group_strip(left, width, color):
        h = Inches(3.65)
        y = Inches(2.2)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(y), int(width), int(h))
        # very dark fill (1/8 of the color)
        r = max(0, color[0] // 9)
        g = max(0, color[1] // 9)
        b = max(0, color[2] // 9)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.color.rgb = RGBColor(color[0] // 4, color[1] // 4, color[2] // 4)
        rect.line.width = Pt(0.5)
        rect.shadow.inherit = False
        # send to back
        spTree = rect._element.getparent()
        spTree.remove(rect._element)
        spTree.insert(2, rect._element)
        return rect

    # Compute strip extents
    def phase_left(i):
        return start_x + i * (box_w + arrow_w)

    pad = Inches(0.18)
    # Group A: phases 0..3 (Deterministic metrics)
    a_left = phase_left(0) - pad
    a_right = phase_left(3) + box_w + pad
    group_strip(a_left, a_right - a_left, RGBColor(0x2D, 0xD4, 0xBF))
    # Group B: phase 4 (Deterministic pre-pass)
    b_left = phase_left(4) - pad
    b_right = phase_left(4) + box_w + pad
    group_strip(b_left, b_right - b_left, RGBColor(0x2D, 0xD4, 0xBF))
    # Group C: phase 5 (AI supervisor)
    c_left = phase_left(5) - pad
    c_right = phase_left(5) + box_w + pad
    group_strip(c_left, c_right - c_left, RGBColor(0xFB, 0x71, 0x85))
    # Group D: phase 6 (Output)
    d_left = phase_left(6) - pad
    d_right = phase_left(6) + box_w + pad
    group_strip(d_left, d_right - d_left, RGBColor(0x81, 0x8C, 0xF8))

    # Group labels above each strip
    def group_label(left, width, text, color):
        y = Inches(1.85)
        tb = slide.shapes.add_textbox(int(left), int(y), int(width), Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = text
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = color
        para.font.name = "Segoe UI"
        para.alignment = PP_ALIGN.CENTER

    group_label(a_left, a_right - a_left, "DETERMINISTIC  ·  metrics", RGBColor(0x2D, 0xD4, 0xBF))
    group_label(b_left, b_right - b_left, "DETERMINISTIC  ·  pre-pass", RGBColor(0x2D, 0xD4, 0xBF))
    group_label(c_left, c_right - c_left, "AI", RGBColor(0xFB, 0x71, 0x85))
    group_label(d_left, d_right - d_left, "OUTPUT", RGBColor(0x81, 0x8C, 0xF8))

    for i, phase in enumerate(PHASES):
        x = phase_left(i)

        # Darken the phase color for box background
        pc = phase["color"]
        r = max(0, pc[0] // 5)
        g = max(0, pc[1] // 5)
        b = max(0, pc[2] // 5)
        box_bg = RGBColor(r, g, b)

        # Box
        box = add_rounded_box(slide, int(x), int(box_y), int(box_w), int(box_h), box_bg)
        box.line.color.rgb = phase["color"]
        box.line.width = Pt(1.5)

        # Number circle
        circle_size = Inches(0.55)
        cx = int(x + (box_w - circle_size) / 2)
        cy = int(box_y + Inches(0.3))
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, int(circle_size), int(circle_size))
        circ.fill.solid()
        circ.fill.fore_color.rgb = phase["color"]
        circ.line.fill.background()
        circ.shadow.inherit = False
        ctf = circ.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.text = phase["num"]
        cp.font.size = Pt(20)
        cp.font.bold = True
        # Phases with light-ish accent colors get dark numerals
        light_accent = i in (2, 3, 4)
        cp.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A) if light_accent else WHITE
        cp.font.name = "Segoe UI"
        cp.alignment = PP_ALIGN.CENTER
        ctf.paragraphs[0].space_before = Pt(0)
        ctf.paragraphs[0].space_after = Pt(0)

        # Title text inside box
        title_y = int(box_y + Inches(1.05))
        ttx = slide.shapes.add_textbox(int(x + Inches(0.08)), title_y, int(box_w - Inches(0.16)), Inches(0.65))
        ttf = ttx.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = phase["title"]
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        tp.font.name = "Segoe UI"
        tp.alignment = PP_ALIGN.CENTER

        # Subtitle text inside box
        sub_y = int(box_y + Inches(1.7))
        stx = slide.shapes.add_textbox(int(x + Inches(0.08)), sub_y, int(box_w - Inches(0.16)), Inches(1.4))
        stf = stx.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = phase["sub"]
        sp.font.size = Pt(11)
        sp.font.color.rgb = MUTED
        sp.font.name = "Segoe UI"
        sp.alignment = PP_ALIGN.CENTER
        sp.line_spacing = Pt(15)

        # Arrow between boxes (not after last)
        if i < n - 1:
            ax = int(x + box_w + Inches(0.02))
            ay = int(box_y + (box_h - arrow_h) / 2)
            add_arrow(slide, ax, ay, int(arrow_w - Inches(0.04)), int(arrow_h))

    # Footer tagline
    foot = slide.shapes.add_textbox(Inches(0), Inches(6.5), prs.slide_width, Inches(0.3))
    fp = foot.text_frame.paragraphs[0]
    fp.text = ("BEFORE score (1–4)  →  Pre-pass guarantees the drop (5)  →  "
               "AI rebalances (6)  →  Recalculate + persist (7)")
    fp.font.size = Pt(12)
    fp.font.color.rgb = ARROW_CLR
    fp.font.name = "Segoe UI"
    fp.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "This slide shows the seven-stage pipeline in execution order. "
        "Stages 1–4 (deterministic metrics) and stage 5 (deterministic pre-pass) require "
        "no LLM. Stage 6 is the only AI step. Stage 7 closes the loop.\n\n"
        "Stage 1 — Ingestion: GitHub issues are pulled into a versioned in-memory IssueCache "
        "via authenticated MCP sync, the rate-limited /demo/api/sync public endpoint, or the "
        "/demo/api/seed test endpoint.\n\n"
        "Stage 2 — Chaos Metrics: ChaosMetricsService evaluates five binary signals first — "
        "mystery meat (≥3 empty bodies), unresolved urgents (≥3), context switching (≥6 in 60 min), "
        "after-hours activity, label sprawl (≥12 distinct labels) — producing a chaos score from 0 to 10. "
        "Runs BEFORE classification.\n\n"
        "Stage 3 — Classification & Compliance: IssueClassifierService sorts each issue into one "
        "of four buckets (deep work, quick win, maintenance, deferred) using first-match-wins label "
        "rules. ComplianceService checks 8 violation rules (score 100 → 0). The 3-3-3 day plan is "
        "built here as well.\n\n"
        "Stage 4 — Stress Score (BEFORE): WorldState aggregates chaos + classification + compliance "
        "into 12 capped variables. Six components are summed: Workload (0–40), Chaos (0–30), "
        "Context Switching (0–15), Clarity (0–10), Sustained (0–15), After Hours (0–10). Total "
        "capped at 100. This is the BEFORE score the demo starts at (~58).\n\n"
        "Stage 5 — Deterministic Pre-pass (reshape only, no LLM): "
        "BurnoutSupervisorService runs two deterministic methods BEFORE any agent is invoked:\n"
        "  • triageUrgent(n) is called directly for every unassigned-urgent issue, stripping the "
        "    `urgent` / `priority:critical` / `priority:high` labels and adding `triaged,backlog`.\n"
        "  • defuseChaosInputs(clock) replaces empty bodies with a scope-pending placeholder and "
        "    normalises after-hours / touched-today timestamps to 10:00 in the demo clock zone.\n"
        "These two calls guarantee the chaos score drops on every reshape, regardless of which "
        "sub-agents the LLM picks. Without this stage the demo's 58 → 8 drop would not be reliable.\n\n"
        "Stage 6 — Supervisor + 6 Sub-Agents (reshape only, AI): "
        "AgenticServices.supervisorBuilder() with maxAgentsInvocations=15 and "
        "SupervisorResponseStrategy.SUMMARY orchestrates six sub-agents — TriageAgent, DeferAgent, "
        "DelegateAgent, ClassifyAgent, ScopeAgent, WellnessAgent — each backed by @Tool methods on "
        "BurnoutMutationTool that emit sealed GitHubAction records (AddLabels, RemoveLabels, Comment, "
        "Unassign, SetBody, SetUpdatedAt). Every agent has a deterministic fallback when the LLM is "
        "unavailable.\n\n"
        "Stage 7 — Output: Mutations are applied to the in-memory IssueCache (not to GitHub). "
        "The ENTIRE deterministic pipeline (stages 2–4) reruns on the mutated issues to produce "
        "the AFTER score (~8 in the demo). A Friday deploy score assesses release readiness "
        "(READY ≥ 80, CAUTION ≥ 50, NOT_READY < 50). Every check-in is persisted as a JPA "
        "StressSnapshot for the longitudinal Study Dashboard."
    )

    out = "docs/Burnout-Pipeline-Sequence.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
