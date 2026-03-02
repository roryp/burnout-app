"""
Generate a widescreen (16:9) PowerPoint presentation for Burnout-as-a-Service.
Uses content from README.md and docs/PSYCHOLOGY.md with images from docs/images/.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
IMG_DIR = os.path.join(os.path.dirname(__file__), "images")
DOCS_DIR = os.path.dirname(__file__)
OUT = os.path.join(os.path.dirname(__file__), "Burnout-as-a-Service.pptx")

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # dark navy
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)    # Microsoft blue
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # green
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)     # red
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # amber
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE_BG = RGBColor(0xF1, 0xF5, 0xF9)      # light slate bg for content slides
CARD_BG = RGBColor(0xE2, 0xE8, 0xF0)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape (for cards/panels)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if alpha is not None:
        from pptx.oxml.ns import qn
        solidFill = shape.fill._fill
        srgbClr = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha_el = srgbClr.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))})
            srgbClr.append(alpha_el)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(tf, items, font_size=16, color=DARK_TEXT, spacing=Pt(6)):
    """Add bullet items to an existing text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0


def add_table(slide, left, top, width, rows_data, col_widths=None, header_color=ACCENT_BLUE,
              font_size=12):
    """Add a formatted table to the slide."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Segoe UI"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else SUBTLE_BG

    return table_shape


def img_path(name):
    """Get image path, checking both docs/images/ and docs/."""
    p = os.path.join(IMG_DIR, name)
    if os.path.exists(p):
        return p
    p2 = os.path.join(DOCS_DIR, name)
    if os.path.exists(p2):
        return p2
    return None


# ═══════════════════════════════════════════════════════════════════════
# BUILD THE PRESENTATION
# ═══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ─── SLIDE 1: Title Slide ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Burnout-as-a-Service", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "AI-Powered Burnout Prevention for Developers", font_size=28,
             color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

# Description
add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(1.2),
             "LangChain4j Supervisor Pattern  •  Azure OpenAI (gpt-5.2)  •  GitHub MCP Integration\n"
             "Analyzes issues, detects stress signals, and organizes your day using the 3-3-3 structure",
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# URL
add_text_box(slide, Inches(3), Inches(6.0), Inches(7), Inches(0.5),
             "Live Demo: https://aka.ms/burnout-app", font_size=18,
             color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 2: The Problem ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Scene 1: The Breaking Point", font_size=36, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(1.0),
             "Alex stares at their screen. 12 open issues, 3 critical bugs, Slack piling up, "
             "and it's 7 PM. No plan, no priorities — just an avalanche of work.",
             font_size=16, color=DARK_TEXT)

p = img_path("scene1-overwhelmed-developer.png")
if p:
    slide.shapes.add_picture(p, Inches(0.8), Inches(2.3), Inches(6.5))

# Psychology callout card
add_shape_bg(slide, Inches(7.8), Inches(2.3), Inches(4.8), Inches(4.2), CARD_BG)
add_text_box(slide, Inches(8.0), Inches(2.5), Inches(4.4), Inches(0.5),
             "Maslach Burnout Inventory", font_size=20, color=ACCENT_RED, bold=True)

tb = slide.shapes.add_textbox(Inches(8.0), Inches(3.2), Inches(4.4), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "Emotional exhaustion — working late, no recovery",
    "Depersonalization — issues blur together",
    "Reduced accomplishment — nothing feels 'done'",
    "",
    "12 context switches have depleted working memory.",
    "Cognitive load is at maximum.",
], font_size=14, color=DARK_TEXT)


# ─── SLIDE 3: The 3-3-3 Structure ─────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Scene 2: The 3-3-3 Day Structure", font_size=36, color=DARK_TEXT, bold=True)

p = img_path("scene2-333-day-structure.png")
if p:
    slide.shapes.add_picture(p, Inches(0.5), Inches(1.5), Inches(7.0))

# Table
data = [
    ["Slot", "Count", "Purpose", "Psychology"],
    ["Deep Work", "1", "Cognitively demanding", "Flow state (Csikszentmihalyi)"],
    ["Quick Wins", "3", "Small, completable", "Dopamine from completion"],
    ["Maintenance", "3", "Routine upkeep", "Low overhead, batch-processable"],
    ["Total", "7", "", "Miller's Law: 7 ± 2"],
]
add_table(slide, Inches(7.8), Inches(1.8), Inches(5.0), data,
          col_widths=[Inches(1.2), Inches(0.7), Inches(1.5), Inches(1.6)], font_size=11)

add_text_box(slide, Inches(7.8), Inches(4.5), Inches(5.0), Inches(0.8),
             "Overflow beyond 7 active items → automatically deferred.\n"
             "Deep work gets a protected 90-minute block.",
             font_size=14, color=MED_GRAY)


# ─── SLIDE 4: Issue Classification ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Scene 3: Issue Classification", font_size=36, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "Deterministic, priority-ordered first-match cascade. No LLM required — pure pattern matching.",
             font_size=16, color=DARK_TEXT)

p = img_path("scene3-issue-classification.png")
if p:
    slide.shapes.add_picture(p, Inches(0.5), Inches(2.0), Inches(7.0))

data = [
    ["Priority", "Bucket", "Label Triggers"],
    ["1st", "DEEP_WORK", "priority:critical, architecture, security, deep-work, performance, RFC"],
    ["2nd", "QUICK_WIN", "good-first-issue, quick-win, size:S, typo, chore, CSS"],
    ["3rd", "MAINTENANCE", "dependencies, documentation, maintenance, tech-debt, refactor"],
    ["4th", "DEFERRED", "Everything else (no matching labels)"],
]
add_table(slide, Inches(7.8), Inches(2.2), Inches(5.0), data,
          col_widths=[Inches(0.7), Inches(1.3), Inches(3.0)], font_size=11)


# ─── SLIDE 5: Stress Score ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Scene 4: Stress Score Algorithm", font_size=36, color=DARK_TEXT, bold=True)

p = img_path("scene4-stress-score.png")
if p:
    slide.shapes.add_picture(p, Inches(0.3), Inches(1.5), Inches(7.2))

data = [
    ["Dimension", "Max", "Formula"],
    ["Workload", "40", "min(20, (assigned−7)×4) + deep work penalties"],
    ["Chaos", "30", "chaosBucket.ordinal() × 10"],
    ["Context Switching", "15", "min(15, (touchedToday−5)×3)"],
    ["Clarity", "15", "min(10, mysteryMeat×2) + min(5, unclearQW)"],
    ["Sustained", "15", "min(15, consecutiveHighDays×5)"],
    ["After-Hours", "10", "min(10, afterHoursIssues×5)"],
]
add_table(slide, Inches(7.8), Inches(1.8), Inches(5.0), data,
          col_widths=[Inches(1.4), Inches(0.6), Inches(3.0)], font_size=11)

# Stress levels
add_shape_bg(slide, Inches(7.8), Inches(5.0), Inches(5.0), Inches(1.8), CARD_BG)
add_text_box(slide, Inches(8.0), Inches(5.1), Inches(4.6), Inches(0.4),
             "Stress Levels", font_size=16, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(8.0), Inches(5.5), Inches(4.6), Inches(1.2),
             "≥ 70  CRITICAL  🔴\n≥ 50  HIGH  🟠\n≥ 30  MODERATE  🟡\n< 30  LOW  🟢",
             font_size=14, color=DARK_TEXT)


# ─── SLIDE 6: Supervisor Agent ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Scene 5: The Supervisor Agent", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "LangChain4j Supervisor Pattern  •  gpt-5.2  •  SUMMARY Strategy  •  Max 10 Invocations",
             font_size=16, color=ACCENT_BLUE)

p = img_path("scene5-supervisor-pattern.png")
if p:
    slide.shapes.add_picture(p, Inches(0.5), Inches(1.7), Inches(12.3))


# ─── SLIDE 7: 5 Sub-Agents Detail ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "5 Sub-Agents + 9 Mutation Tools", font_size=36, color=DARK_TEXT, bold=True)

data = [
    ["Agent", "Role", "Tools", "Psychology"],
    ["DeferAgent", "Push non-critical to next sprint", "deferIssue()", "Reduces cognitive load by shrinking active set"],
    ["DelegateAgent", "Redistribute across team", "delegateIssue()", "\"I don't have to do everything\""],
    ["ClassifyAgent", "Organize into 3-3-3", "markAsDeepWork(), classifyAsQuickWin(), classifyAsMaintenance()", "Imposes order on chaos"],
    ["ScopeAgent", "Flag unclear issues", "addScopeNeeded()", "Eliminates 'mystery meat' ambiguity"],
    ["WellnessAgent", "Recommend self-care", "suggestBreak(), slowIntake(), blockCalendarTime()", "Direct burnout prevention"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), data,
          col_widths=[Inches(1.5), Inches(2.5), Inches(4.3), Inches(4.0)], font_size=12)

# 3 AI personas
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(5.0), Inches(0.5),
             "3 AI Personas", font_size=20, color=DARK_TEXT, bold=True)
data2 = [
    ["Service", "Persona", "Purpose"],
    ["ExplainerAiService", "Supportive productivity coach", "Explain the why behind the plan"],
    ["ProtectiveAiService", "Protective AI companion", "Emotional support (Plutchik model)"],
    ["FridayDeployAiService", "Calm release engineer", "Deploy risk assessment"],
]
add_table(slide, Inches(0.5), Inches(5.5), Inches(8.0), data2,
          col_widths=[Inches(2.5), Inches(2.5), Inches(3.0)], font_size=11)


# ─── SLIDE 8: Algorithm Pipeline ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "Algorithm Pipeline — End to End", font_size=36, color=WHITE, bold=True)

p = img_path("algorithm-pipeline.png")
if p:
    slide.shapes.add_picture(p, Inches(0.3), Inches(1.0), Inches(12.7))


# ─── SLIDE 9: Flamegraph — Before ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Scene 6: The Flamegraph — BEFORE", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "12 issues with no labels → all dumped into Deferred. No structure at all.",
             font_size=18, color=ACCENT_RED)

p = img_path("flamegraph-before.png")
if p:
    slide.shapes.add_picture(p, Inches(0.5), Inches(1.7), Inches(12.3))


# ─── SLIDE 10: Flamegraph — After ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Scene 7: The Flamegraph — AFTER", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Same 12 issues → now structured into 1 Deep Work, 3 Quick Wins, 3 Maintenance, 5 Deferred.",
             font_size=18, color=ACCENT_GREEN)

p = img_path("flamegraph-after.png")
if p:
    slide.shapes.add_picture(p, Inches(0.5), Inches(1.7), Inches(12.3))


# ─── SLIDE 11: Before vs After ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Scene 8: The Transformation", font_size=36, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "Same issues, same deadlines — but structured. That's the difference between burnout and balance.",
             font_size=18, color=DARK_TEXT)

data = [
    ["Metric", "Before (No Labels)", "After (Reshaped)"],
    ["Deep Work", "0", "1"],
    ["Quick Wins", "0", "3"],
    ["Maintenance", "0", "3"],
    ["Deferred", "12 (all)", "5 (intentional)"],
    ["Stress Score", "25/100", "20/100"],
    ["3-3-3 Compliant", "No ❌", "Yes ✅"],
]
add_table(slide, Inches(2.5), Inches(2.2), Inches(8.0), data,
          col_widths=[Inches(2.5), Inches(2.75), Inches(2.75)], font_size=16)


# ─── SLIDE 12: Chaos & Compliance ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Chaos Metrics & Compliance", font_size=36, color=DARK_TEXT, bold=True)

# Chaos table
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.5),
             "Chaos Score (0–10) — 5 Binary Signals × 2 pts", font_size=16, color=ACCENT_RED, bold=True)

data_chaos = [
    ["Signal", "Trigger", "What It Reveals"],
    ["Mystery meat", "≥ 3 blank issues", "Poor issue quality"],
    ["Unresolved urgent", "≥ 3 items > 24h", "Broken priorities"],
    ["Touched today", "≥ 6 in 60 min", "Reactive firefighting"],
    ["After-hours", "Outside 8am–6pm", "Boundary erosion"],
    ["Label explosion", "≥ 12 labels", "Taxonomy chaos"],
]
add_table(slide, Inches(0.3), Inches(1.9), Inches(6.2), data_chaos,
          col_widths=[Inches(1.5), Inches(1.7), Inches(3.0)], font_size=11,
          header_color=ACCENT_RED)

# Compliance table
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(6), Inches(0.5),
             "Compliance (100 → 0) — 8 Violations", font_size=16, color=ACCENT_AMBER, bold=True)

data_comp = [
    ["Violation", "Severity", "Deduction"],
    ["MULTIPLE_DEEP_WORK", "CRITICAL", "−25"],
    ["TOTAL_OVERLOAD", "CRITICAL", "−25"],
    ["CONTEXT_SWITCHING", "CRITICAL", "−25"],
    ["QUICK_WIN_OVERLOAD", "WARNING", "−10"],
    ["MAINTENANCE_OVERLOAD", "WARNING", "−10"],
    ["UNCLEAR_QUICK_WINS", "WARNING", "−10"],
    ["NO_DEEP_WORK", "INFO", "−5"],
    ["DEFERRED_GROWING", "INFO", "−5"],
]
add_table(slide, Inches(6.8), Inches(1.9), Inches(6.0), data_comp,
          col_widths=[Inches(2.8), Inches(1.2), Inches(2.0)], font_size=11,
          header_color=ACCENT_AMBER)


# ─── SLIDE 13: Emotional Detection & Protection ──────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Emotional Detection & Protective Intervention", font_size=36, color=DARK_TEXT, bold=True)

# Plutchik table
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
             "Plutchik Model — 4 Emotions Detected", font_size=16, color=ACCENT_BLUE, bold=True)

data_emo = [
    ["Emotion", "Family", "Signals"],
    ["Frustration", "Anger", "Context switches > 5, blocked items"],
    ["Exhaustion", "Sadness", "After-hours, consecutive high days ≥ 2"],
    ["Overwhelm", "Surprise", "Deep work > 1, assigned > 10"],
    ["Anxiety", "Fear", "Stale urgent items, mystery meat"],
]
add_table(slide, Inches(0.3), Inches(1.9), Inches(6.2), data_emo,
          col_widths=[Inches(1.3), Inches(1.0), Inches(3.9)], font_size=12)

# Protection triggers
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(6), Inches(0.5),
             "Protection Circuit Breaker", font_size=16, color=ACCENT_RED, bold=True)

data_prot = [
    ["Trigger", "Threshold"],
    ["Sustained stress", "consecutiveHighDays ≥ 2"],
    ["Boundary erosion", "hasAfterHoursActivity()"],
    ["Acute overload", "stressScore ≥ 70"],
    ["Cognitive capacity", "totalAssigned > 10"],
]
add_table(slide, Inches(6.8), Inches(1.9), Inches(6.0), data_prot,
          col_widths=[Inches(2.5), Inches(3.5)], font_size=12,
          header_color=ACCENT_RED)

# AI response principles
add_shape_bg(slide, Inches(6.8), Inches(4.0), Inches(6.0), Inches(2.8), CARD_BG)
add_text_box(slide, Inches(7.0), Inches(4.1), Inches(5.6), Inches(0.4),
             "AI Response Principles", font_size=16, color=DARK_TEXT, bold=True)
tb = slide.shapes.add_textbox(Inches(7.0), Inches(4.6), Inches(5.6), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "✓ Validate without patronizing",
    "✓ Concrete suggestions only",
    "✓ Brevity — stressed people have reduced reading comprehension",
    "✓ No guilt or shame — focus on self-care",
    "✓ One actionable item — decision fatigue is real",
], font_size=13, color=DARK_TEXT)


# ─── SLIDE 14: Flamegraph Psychology ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Flamegraph Visualization Psychology", font_size=36, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "The flamegraph is a stress communication tool — fire metaphors trigger the brain's threat detection system.",
             font_size=18, color=DARK_TEXT)

# Design elements
data_vis = [
    ["Element", "Psychology", "Why It Works"],
    ["🔥 Fire = Danger", "Primal association", "Immediate emotional response"],
    ["📊 Height = Depth", "Taller bars feel 'heavier'", "Visual weight = cognitive weight"],
    ["🚦 Color = Urgency", "Red / Amber / Green", "Traffic light intuition (universal)"],
    ["↔️ Width = Proportion", "Wider = more impact", "Relative comparison at a glance"],
]
add_table(slide, Inches(0.5), Inches(2.0), Inches(7.0), data_vis,
          col_widths=[Inches(1.8), Inches(2.5), Inches(2.7)], font_size=13)

# Per-issue stress
add_shape_bg(slide, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.5), CARD_BG)
add_text_box(slide, Inches(8.2), Inches(2.1), Inches(4.4), Inches(0.4),
             "Per-Issue Stress Formula", font_size=16, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(8.2), Inches(2.6), Inches(4.4), Inches(0.5),
             "stress = base + (complexity × 3) + (global × 0.3) + bonuses",
             font_size=13, color=ACCENT_BLUE, bold=True)

data_base = [
    ["Category", "Base", "Bonuses"],
    ["Deep Work", "60", "urgent/critical: +20"],
    ["Maintenance", "30", "bug: +10"],
    ["Quick Wins", "20", "—"],
    ["Deferred", "10", "—"],
]
add_table(slide, Inches(8.2), Inches(3.3), Inches(4.4), data_base,
          col_widths=[Inches(1.4), Inches(0.8), Inches(2.2)], font_size=12)

add_text_box(slide, Inches(8.2), Inches(5.5), Inches(4.4), Inches(0.8),
             "Colors: < 35% 🟢  •  35–64% 🟡  •  ≥ 65% 🔴\n"
             "Global stress leak (×0.3): high-stress environments make every task feel harder.",
             font_size=12, color=MED_GRAY)


# ─── SLIDE 15: Friday Deploy & Calendar ───────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Friday Deploy & Calendar Fragmentation", font_size=36, color=DARK_TEXT, bold=True)

# Friday deploy
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
             "Friday Deploy Score (0–100)", font_size=18, color=ACCENT_RED, bold=True)

data_fri = [
    ["Condition", "Deduction"],
    ["chaos > 5", "−20"],
    ["chaos > 8", "−20 (cumulative: −40)"],
    ["!isCompliant", "−15"],
    ["urgentUnassigned > 0", "−15"],
    ["afterHoursSignal", "−10"],
    ["mysteryMeat > 3", "−10"],
]
add_table(slide, Inches(0.3), Inches(1.9), Inches(5.5), data_fri,
          col_widths=[Inches(3.0), Inches(2.5)], font_size=13,
          header_color=ACCENT_RED)

add_text_box(slide, Inches(0.3), Inches(5.2), Inches(5.5), Inches(1.5),
             "≥ 80  READY 🟢  •  50–79  CAUTION 🟡  •  < 50  NOT READY 🔴\n\n"
             "Counters optimism bias and completion bias:\n"
             "\"It'll be fine\" + \"Just finish it before the weekend\"",
             font_size=14, color=DARK_TEXT)

# Calendar Fragmentation
add_shape_bg(slide, Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(6.7), Inches(1.5), Inches(5.9), Inches(0.5),
             "Calendar Fragmentation", font_size=18, color=ACCENT_BLUE, bold=True)

tb = slide.shapes.add_textbox(Inches(6.7), Inches(2.2), Inches(5.9), Inches(4.0))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "Deep Work Feasibility Threshold:",
    "  largestFreeBlock ≥ 90 minutes",
    "",
    "Why 90 minutes?",
    "  • 23 min — ramp-up to flow state",
    "  • 60 min — productive deep work",
    "  • 7 min — buffer",
    "",
    "No 90-min block → calendarBlocked = true",
    "→ deep work item deferred automatically",
], font_size=14, color=DARK_TEXT)


# ─── SLIDE 16: Graceful Degradation ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Graceful Degradation", font_size=36, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "Core principle: every AI feature must work without AI.",
             font_size=22, color=ACCENT_BLUE, bold=True)

data_deg = [
    ["Level", "State", "Behavior"],
    ["0", "Full LLM available", "AI explanations, emotional support, supervisor active"],
    ["1", "LLM call fails", "Catch exception → fallback; all metrics still accurate"],
    ["2", "LLM not configured", "llmEnabled = false → full deterministic with pre-written messages"],
    ["3", "Backend partial", "Health returns UP; endpoints degrade independently"],
]
add_table(slide, Inches(1.0), Inches(2.2), Inches(11.0), data_deg,
          col_widths=[Inches(0.8), Inches(2.2), Inches(8.0)], font_size=14)

# Why it matters
add_shape_bg(slide, Inches(1.0), Inches(4.8), Inches(11.0), Inches(2.0), CARD_BG)
add_text_box(slide, Inches(1.2), Inches(4.9), Inches(10.6), Inches(0.4),
             "Why This Matters", font_size=18, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.4), Inches(10.6), Inches(1.2),
             "If a burnout prevention tool fails when you need it most, it increases stress instead of reducing it. "
             "All metrics are deterministic — stress scores, chaos, compliance, classification, and day plans "
             "compute without any LLM. AI agents only add the how (mutation plan) and why (explanation).",
             font_size=15, color=DARK_TEXT)


# ─── SLIDE 17: Theoretical Foundations ────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Theoretical Foundations", font_size=36, color=DARK_TEXT, bold=True)

data_theory = [
    ["Framework", "Key Concept", "Implementation"],
    ["Maslach Burnout Inventory", "3 dimensions: exhaustion,\ndepersonalization, accomplishment", "After-hours → exhaustion\nBlank issues → depersonalization\nNo deep work → reduced accomplishment"],
    ["Cognitive Load Theory\n(Sweller, 1988)", "Working memory is limited;\nextraneous load = enemy", "3-3-3 cap at 7 items;\nclassification removes ambiguity"],
    ["Yerkes-Dodson Law (1908)", "Performance peaks at\nmoderate stress", "Targets 30–50 (MODERATE)\nas optimal zone"],
    ["Deep Work (Newport, 2016)", "90+ min uninterrupted\nblocks for flow", "Calendar fragmentation check;\nexactly 1 deep work per day"],
    ["Pomodoro / Time Boxing", "Short intervals with\nbreaks maintain energy", "Quick wins as break-points\nbetween deep work sessions"],
    ["Plutchik's Wheel (2001)", "8 primary emotions with\nbehavioral signatures", "4 emotions detected from\nGitHub signals"],
]
add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), data_theory,
          col_widths=[Inches(2.6), Inches(3.8), Inches(6.3)], font_size=12)


# ─── SLIDE 18: Architecture & Usage ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Architecture & MCP Usage", font_size=36, color=DARK_TEXT, bold=True)

# Agent flow
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
             "End-to-End Flow", font_size=18, color=ACCENT_BLUE, bold=True)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(6.0), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "1. Sync — GitHub issues → IssueCache (ConcurrentHashMap)",
    "2. Calculate — ChaosMetrics, Compliance, Classification",
    "3. Build WorldState — 18 capped variables",
    "4. Invoke Supervisor — planner coordinates 5 sub-agents",
    "5. Accumulate Mutations — @Tool methods → pendingActions",
    "6. Return Response — plan, scores, protective messages",
], font_size=14, color=DARK_TEXT, spacing=Pt(10))

# MCP tools
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(6), Inches(0.5),
             "4 MCP Tools (VS Code Copilot Chat)", font_size=18, color=ACCENT_BLUE, bold=True)

data_mcp = [
    ["Tool", "Description"],
    ["sync_issues", "Fetch GitHub issues and sync to backend"],
    ["show_burnout_wheel", "Interactive flamegraph with 3-3-3 plan"],
    ["reshape_day", "AI-powered optimization — applies labels"],
    ["get_stress_score", "Stress score 0–100 (LOW → CRITICAL)"],
]
add_table(slide, Inches(6.8), Inches(1.9), Inches(6.0), data_mcp,
          col_widths=[Inches(2.2), Inches(3.8)], font_size=13)

# Usage
add_text_box(slide, Inches(7.0), Inches(4.2), Inches(6), Inches(0.5),
             "Usage in VS Code", font_size=18, color=ACCENT_BLUE, bold=True)

add_shape_bg(slide, Inches(6.8), Inches(4.8), Inches(6.0), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B))
add_text_box(slide, Inches(7.0), Inches(4.9), Inches(5.6), Inches(2.0),
             "Sync issues for owner/repo\n"
             "Show my burnout wheel for owner/repo\n"
             "Reshape my day for owner/repo\n"
             "What's my stress score for owner/repo",
             font_size=14, color=ACCENT_GREEN, font_name="Consolas")


# ─── SLIDE 19: Quick Start ───────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SUBTLE_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Quick Start", font_size=36, color=DARK_TEXT, bold=True)

# Azure deploy
add_shape_bg(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
add_text_box(slide, Inches(0.7), Inches(1.4), Inches(5.4), Inches(0.5),
             "Option A: Deploy to Azure", font_size=20, color=ACCENT_BLUE, bold=True)

tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(5.4), Inches(4.2))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "Prerequisites:",
    "  • GitHub CLI — gh auth login",
    "  • VS Code with GitHub Copilot",
    "  • Node.js 18+",
    "",
    "One command:",
    "  azd auth login",
    "  azd up",
    "",
    "Provisions: Container Apps, Azure OpenAI,",
    "  ACR, managed identity (Bicep)",
], font_size=14, color=DARK_TEXT)

# Local
add_shape_bg(slide, Inches(6.7), Inches(1.3), Inches(6.1), Inches(5.5), WHITE)
add_text_box(slide, Inches(6.9), Inches(1.4), Inches(5.7), Inches(0.5),
             "Option B: Run Locally", font_size=20, color=ACCENT_BLUE, bold=True)

tb = slide.shapes.add_textbox(Inches(6.9), Inches(2.1), Inches(5.7), Inches(4.2))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "Requires: Java 21+, Maven",
    "",
    "Backend:",
    "  cd backend",
    "  mvn clean package -DskipTests",
    "  java -jar target/burnout-backend-*.jar",
    "",
    "MCP App:",
    "  cd mcp-app",
    "  npm install && npm run build",
    "",
    "Reload VS Code — .vscode/mcp.json pre-configured",
], font_size=14, color=DARK_TEXT)


# ─── SLIDE 20: Closing Slide ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Burnout-as-a-Service", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
             "Not less work — structured work.", font_size=28,
             color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(1.5),
             "Live Demo: https://aka.ms/burnout-app\n"
             "GitHub: github.com/roryp/burnout-app\n"
             "Psychology Deep-Dive: docs/PSYCHOLOGY.md",
             font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.5),
             "Java 21  •  Spring Boot 3  •  LangChain4j  •  Azure OpenAI  •  MCP  •  GitHub Copilot",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"✅ Presentation saved to: {OUT}")
print(f"   Slides: {len(prs.slides)}")
