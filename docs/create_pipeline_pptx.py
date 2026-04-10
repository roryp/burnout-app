"""Generate a 6-slide PPTX for the Burnout-as-a-Service algorithm pipeline."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

STAGE_COLORS = [
    RGBColor(0x3B, 0x82, 0xF6),  # 1 blue
    RGBColor(0xA7, 0x8B, 0xFA),  # 2 purple
    RGBColor(0xF5, 0x9E, 0x0B),  # 3 amber
    RGBColor(0x2D, 0xD4, 0xBF),  # 4 teal
    RGBColor(0xFB, 0x71, 0x85),  # 5 pink
    RGBColor(0x81, 0x8C, 0xF8),  # 6 indigo
]

SLIDES = [
    {
        "stage": "Stage 1",
        "title": "Ingestion",
        "points": [
            ("GitHub Issues → In-Memory Cache",
             "Issues synced via MCP or public API (1 per repo per 5 min)"),
            ("9 Fields Per Issue",
             "number, title, body, labels, assignees, createdAt, updatedAt, state, milestone"),
            ("Demo Labels Override Timestamps",
             "demo:touched-today, demo:after-hours, demo:stale-14d"),
        ],
        "notes": (
            "The ingestion stage pulls GitHub issues into a versioned in-memory cache. "
            "There are two paths: the MCP tool syncs with a GitHub Bearer token for authenticated access, "
            "while the public demo API is rate-limited to one sync per repo every five minutes to stay within GitHub's unauthenticated limits.\n\n"
            "Each issue is stored as a Java record with nine fields — the key ones for stress analysis are labels, "
            "assignees, and the camelCase timestamps createdAt and updatedAt. Using snake_case here silently produces nulls, "
            "which is a common gotcha.\n\n"
            "For demo scenarios, synthetic demo labels like demo:touched-today override real timestamps entirely. "
            "This is the golden rule: if any demo label is present on an issue, the system never consults the actual timestamp fields. "
            "This lets us reliably reproduce stress scenarios regardless of when the demo is run."
        ),
    },
    {
        "stage": "Stage 2",
        "title": "Chaos Metrics",
        "points": [
            ("Chaos Score 0–10 (5 signals × 2 pts)",
             "mystery meat, urgent, context switching, after-hours, label sprawl"),
            ("Runs BEFORE Classification",
             "ChaosMetricsService.calculate() is the first call after loading issues"),
            ("Timezone-Aware After-Hours Detection",
             "Working hours: 9 AM – 6 PM; weekends always count"),
        ],
        "notes": (
            "Chaos metrics run first — before classification and before compliance. "
            "The chaos score is a simple additive metric from zero to ten. "
            "Five boolean signals each contribute two points: three or more mystery-meat issues without descriptions, "
            "three or more unresolved urgents older than 24 hours, six or more issues touched in the last 60 minutes, "
            "any after-hours activity, and twelve or more distinct labels across all issues indicating label sprawl.\n\n"
            "This runs before classification because the chaos score feeds into the WorldState, "
            "which needs both chaos AND classification data to compute the stress score.\n\n"
            "After-hours detection is fully timezone-aware. The checkin page sends the browser's timezone, "
            "and the server uses it to determine if updates happened before 9 AM, after 6 PM, or on weekends. "
            "For demo scenarios, the demo:after-hours label overrides real timestamps entirely."
        ),
    },
    {
        "stage": "Stage 3",
        "title": "Classification & Compliance",
        "points": [
            ("4 Buckets — First Match Wins",
             "DEEP_WORK → QUICK_WIN → MAINTENANCE → DEFERRED"),
            ("Compliance Score 100 → 0 (8 violations)",
             "CRITICAL −25 · WARNING −10 · INFO −5"),
            ("3-3-3 Day Structure (classification runs twice)",
             "1 deep work + 3 quick wins + 3 maintenance = 7 active max"),
        ],
        "notes": (
            "Classification runs inside ComplianceService.analyze(), which is the second call after chaos metrics. "
            "Every issue is sorted into exactly one of four buckets using a first-match-wins strategy. "
            "Deep work catches anything with priority:critical, security, architecture, or issues estimated over two hours. "
            "Quick wins match good-first-issue, trivial, or short tasks with clear scope. "
            "Maintenance covers dependencies, documentation, chores, and CI. Everything else falls into deferred.\n\n"
            "Compliance starts at 100 and deducts points for eight violation types. "
            "Critical violations like having multiple deep work items or more than seven active issues cost 25 points each. "
            "Warnings like bucket overflow cost 10, and informational issues like a growing deferred backlog cost 5.\n\n"
            "The 3-3-3 day structure is the core framework: at most one deep work item, three quick wins, "
            "and three maintenance tasks — seven active issues maximum. "
            "Notably, classification runs a second time inside buildDayPlan() to arrange the sorted issues. "
            "Anything beyond the seven active slots is automatically deferred."
        ),
    },
    {
        "stage": "Stage 4",
        "title": "Stress Score (0–100)",
        "points": [
            ("6-Component Breakdown",
             "Workload · Chaos · Context Switching · Clarity · Sustained · After Hours"),
            ("4 Stress Levels",
             "≥ 70 CRITICAL · ≥ 50 HIGH · ≥ 30 MODERATE · < 30 LOW"),
            ("12 Capped World State Variables",
             "All inputs capped to prevent overflow and keep scores deterministic"),
        ],
        "notes": (
            "The stress score is the sum of six independently calculated components, each with its own cap. "
            "Workload goes up to 40 — it penalises having more than seven assigned issues and having multiple deep work items. "
            "Chaos maps the bucket level to 0, 10, 20, or 30 points. "
            "Context switching adds three points for each issue touched beyond the first five, up to 15. "
            "Clarity adds two points per mystery-meat issue up to 10. "
            "Sustained load adds five points per consecutive high-chaos day up to 15. "
            "After-hours adds five points per issue updated outside working hours, capped at 10.\n\n"
            "The total is clamped to 100 and mapped to four stress levels: "
            "70 and above is critical, 50 is high, 30 is moderate, and below 30 is low.\n\n"
            "All twelve input variables in the world state are individually capped — for example, "
            "total assigned issues cap at 15 and deferred at 10. "
            "This makes the scoring fully deterministic and prevents any single runaway metric from dominating."
        ),
    },
    {
        "stage": "Stage 5",
        "title": "AI Agents",
        "points": [
            ("LangChain4j Supervisor — 5 Sub-Agents",
             "Defer · Delegate · Classify · Scope · Wellness"),
            ("Protection Triggers Automatically",
             "2+ high-chaos days, after-hours, stress ≥ 70, or 10+ issues"),
            ("Deterministic Fallback — Always Works",
             "Every agent has a fallback when the LLM is unavailable"),
        ],
        "notes": (
            "The AI layer uses LangChain4j's supervisor pattern with Azure OpenAI gpt-4o at temperature 0.3. "
            "Five specialised sub-agents each have tool methods that apply GitHub labels: "
            "the defer agent marks issues for next sprint, the delegate agent flags items for reassignment, "
            "the classify agent sorts issues into the 3-3-3 buckets, the scope agent tags unclear issues as blocked, "
            "and the wellness agent can suggest breaks, slow intake, or block calendar time.\n\n"
            "Protection kicks in automatically when the system detects sustained risk: "
            "two or more consecutive high-chaos days, any after-hours activity, a stress score of 70 or above, "
            "or more than ten assigned issues. When triggered, the protective AI service generates "
            "a supportive intervention message alongside the regular analysis.\n\n"
            "Every AI feature has a deterministic fallback path. If the LLM is unavailable — "
            "whether from dummy credentials, network issues, or expired tokens — "
            "all features continue to work with pre-computed responses. "
            "The AI explains; it never decides. All metrics are calculated deterministically first."
        ),
    },
    {
        "stage": "Stage 6",
        "title": "Output & Persistence",
        "points": [
            ("GitHub Mutation Plan",
             "Labels and comments applied to reshape your day into 3-3-3"),
            ("Friday Deploy Score (0–100)",
             "READY ≥ 80 · CAUTION ≥ 50 · NOT_READY < 50"),
            ("Stress Snapshots Persisted (JPA)",
             "All 6 breakdown components + self-report saved for trend analysis"),
        ],
        "notes": (
            "The mutation plan is a sealed interface with three action types: add labels, remove labels, and comment. "
            "The supervisor agent's tool calls are collected into this plan, which the MCP app then executes against GitHub. "
            "Labels like deferred, quick-win, deep-work, 3-3-3, and needs-scope are applied to reshape "
            "the developer's workload into the target structure.\n\n"
            "The Friday deploy score starts at 100 and deducts points for risk factors: "
            "20 points if chaos exceeds 5, another 20 if it exceeds 8, 15 for unresolved urgents, "
            "15 for non-compliance, 10 for after-hours activity, and 10 for more than three mystery-meat issues. "
            "The result maps to three statuses: ready at 80 or above, caution from 50 to 79, and not ready below 50.\n\n"
            "Every stress check-in is persisted as a JPA entity with all six breakdown components, "
            "the overall score, stress level, and an optional self-reported score and note. "
            "This powers the study dashboard's trend charts, letting researchers track "
            "how individual participants' stress evolves over time across multiple check-ins."
        ),
    },
]


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide(prs, data, idx):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, BG)

    stage_color = STAGE_COLORS[idx]
    sw = prs.slide_width
    margin = Inches(1.2)
    content_w = sw - margin * 2

    # Stage label
    tx = slide.shapes.add_textbox(margin, Inches(0.55), content_w, Inches(0.35))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["stage"].upper()
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = stage_color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(0)

    # Title
    tx2 = slide.shapes.add_textbox(margin, Inches(0.9), content_w, Inches(0.65))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = data["title"]
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(0)

    # 3 points — single textbox for all, using paragraphs
    left = Inches(1.6)
    top = Inches(2.0)
    box_w = sw - Inches(3.2)
    box_h = Inches(3.2)
    tx_pts = slide.shapes.add_textbox(left, top, box_w, box_h)
    tf_pts = tx_pts.text_frame
    tf_pts.word_wrap = True

    for i, (heading, detail) in enumerate(data["points"]):
        # Heading paragraph
        if i == 0:
            p_h = tf_pts.paragraphs[0]
        else:
            p_h = tf_pts.add_paragraph()
        p_h.text = heading
        p_h.font.size = Pt(20)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.font.name = "Segoe UI"
        p_h.space_before = Pt(18) if i > 0 else Pt(0)
        p_h.space_after = Pt(2)

        # Detail paragraph
        p_d = tf_pts.add_paragraph()
        p_d.text = detail
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = MUTED
        p_d.font.name = "Segoe UI"
        p_d.space_after = Pt(0)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = data["notes"]


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)  # 16:9 widescreen

    for idx, data in enumerate(SLIDES):
        add_slide(prs, data, idx)

    out = "docs/Burnout-Pipeline.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
