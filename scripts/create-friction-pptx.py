"""Create a PPTX with the 7 friction point images"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

IMG_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "images", "friction")
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "friction-points.pptx")

SLIDES = [
    {
        "title": "The 7 Friction Points Causing Developer Burnout",
        "subtitle": "How Burnout-as-a-Service detects and mitigates each one",
        "image": None,
    },
    {
        "title": "1. Fragmentation",
        "subtitle": "Too many tools, tabs, models, prompts, and context switches",
        "image": "01-fragmentation.png",
    },
    {
        "title": "2. Learning Pressure",
        "subtitle": "The constant feeling that you're already behind",
        "image": "02-learning-pressure.png",
    },
    {
        "title": "3. Performance Standards",
        "subtitle": "Pressure to produce more, faster, because AI exists",
        "image": "03-performance-standards.png",
    },
    {
        "title": "4. Isolation",
        "subtitle": "Doing cognitively heavy work alone, with less human grounding",
        "image": "04-isolation.png",
    },
    {
        "title": "5. Interface Friction",
        "subtitle": "Fighting clunky tools instead of flowing through the work",
        "image": "05-interface-friction.png",
    },
    {
        "title": "6. Altitude Sickness",
        "subtitle": "Strategic thinking pulled into fine-grained review and execution",
        "image": "06-altitude-sickness.png",
    },
    {
        "title": "7. Workload Creep",
        "subtitle": "AI makes more possible, so more gets expected",
        "image": "07-workload-creep.png",
    },
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Title slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = SLIDES[0]["title"]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9.333), Inches(1))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = SLIDES[0]["subtitle"]
p2.font.size = Pt(22)
p2.font.color.rgb = RGBColor(0x88, 0xCC, 0xFF)
p2.alignment = PP_ALIGN.CENTER

# --- Image slides ---
for s in SLIDES[1:]:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)

    # Title at top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s["title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(0.9), Inches(11.333), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = s["subtitle"]
    p2.font.size = Pt(18)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x88, 0xCC, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

    # Image — centered, large
    img_path = os.path.join(IMG_DIR, s["image"])
    if os.path.exists(img_path):
        # Image is 1536x1024 (3:2). Fit within remaining slide area.
        img_w = Inches(10.5)
        img_h = Inches(5.6)  # slightly under to leave margin at bottom
        left = (prs.slide_width - img_w) // 2
        top = Inches(1.55)
        slide.shapes.add_picture(img_path, left, top, img_w, img_h)

prs.save(OUT)
print(f"PPTX saved to {OUT}")
